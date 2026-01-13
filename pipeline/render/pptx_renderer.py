from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from pipeline.contracts.contract_io import load_contract


def render_pptx_from_contract(
    contract_path: Path,
    template_path: Path,
    output_path: Path
) -> None:
    """
    Render PPTX usando UNA plantilla base con layouts reales.
    """

    contract = load_contract(contract_path)
    slides_def = contract.get("slides", [])

    #  Cargar plantilla base
    presentation = Presentation(template_path)

    for slide_def in slides_def:
        slide_type = slide_def["slide_type"]

        # ───────── TITLE SLIDE ─────────
        if slide_type == "title":
            layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(layout)
            _set_title(slide, slide_def["title_text"])

        # ───────── CONTENT SLIDE ─────────
        elif slide_type == "content":
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            _set_title(slide, slide_def["title_text"])

            content = slide_def["content"]

            if content["content_type"] == "paragraph":
                _set_body_paragraph(slide, content["text"])

            elif content["content_type"] == "bullets":
                _set_body_bullets(slide, content["bullets"])

    presentation.save(output_path)


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def _set_title(slide, text: str) -> None:
    """
    Inserta el texto en el placeholder de título.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            shape.text = text
            return

    print("[WARN] TITLE placeholder no encontrado.")


def _set_body_paragraph(slide, text: str) -> None:
    """
    Inserta texto como párrafo único.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type != PP_PLACEHOLDER.TITLE:
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.text = text
            return

    print("[WARN] No se encontró placeholder de contenido (paragraph).")


def _set_body_bullets(slide, bullets) -> None:
    """
    Inserta texto como bullets reales de PowerPoint.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type != PP_PLACEHOLDER.TITLE:
            text_frame = shape.text_frame
            text_frame.clear()

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet
                p.level = 0  # nivel de bullet

            return

    print("[WARN] No se encontró placeholder de contenido (bullets).")